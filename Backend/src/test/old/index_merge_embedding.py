import PyPDF2
import sys
from os import listdir
from os.path import isfile, join, isdir, abspath, basename, relpath
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import TokenTextSplitter
from pptx import Presentation
import docx
from elasticsearch import Elasticsearch
from transformers import AutoModel, AutoTokenizer
import numpy as np

def get_files(dir):
    file_list = []
    for f in listdir(dir):
        if isfile(join(dir, f)):
            file_list.append(join(dir, f))
        elif isdir(join(dir, f)):
            file_list = file_list + get_files(join(dir, f))
    return file_list

def getTextFromWord(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)

def getTextFromPPTX(filename):
    prs = Presentation(filename)
    fullText = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                fullText.append(shape.text)
    return '\n'.join(fullText)

def file_already_indexed(es, index_name, file_name, file_url):
    query = {
        "query": {
            "bool": {
                "must": [
                    {"term": {"metadata.filename.keyword": file_name}},
                    {"term": {"metadata.url.keyword": file_url}}
                ]
            }
        }
    }
    response = es.search(index=index_name, body=query)
    return response['hits']['total']['value'] > 0

def main_indexing(mypath):
    # Local model directory
    model_name = "./local_model"
    model_kwargs = {'device': 'cpu'}
    encode_kwargs = {'normalize_embeddings': True}
    hf = HuggingFaceEmbeddings(
        model_name=model_name,
        model_kwargs=model_kwargs,
        encode_kwargs=encode_kwargs
    )

    # Elasticsearch setup
    es = Elasticsearch("https://192.168.140.243:9200",
                       ssl_assert_fingerprint="5c7375ff8f27f77b34410f674e0b000bdd40d2d688408a4357aee962abe96265",
                       basic_auth=("elastic", "welcome")
    )

    # Base URL for Apache
    base_url = "http://192.168.140.238"

    print("Indexing...")
    onlyfiles = get_files(mypath)
    for file in onlyfiles:
        file_name = basename(file)
        relative_path = relpath(file, mypath)
        
        # Construct the file URL using the relative path
        file_url = f"{base_url}/{relative_path.replace('\\', '/')}"
        print(file_url)

        # Check if the file is already indexed
        if file_already_indexed(es, "search-my_elastic-2", file_name, file_url):
            print(f"File {file_name} is already indexed. Skipping...")
            continue

        file_content = ""
        if file.endswith(".pdf"):
            print("indexing " + file)
            reader = PyPDF2.PdfReader(file)
            for i in range(len(reader.pages)):
                file_content += " " + reader.pages[i].extract_text()
        elif file.endswith((".txt", ".md", ".markdown")):
            print("indexing " + file)
            with open(file, 'r') as f:
                file_content = f.read()
        elif file.endswith(".docx"):
            print("indexing " + file)
            file_content = getTextFromWord(file)
        elif file.endswith(".pptx"):
            print("indexing " + file)
            file_content = getTextFromPPTX(file)
        else:
            continue

        # Split the text and generate embeddings
        #text_splitter = TokenTextSplitter(chunk_size=800, chunk_overlap=80)
        #texts = text_splitter.split_text(file_content)
        #embeddings = hf.embed_documents(texts)
        def custom_text_splitter(text, chunk_size, chunk_overlap):
            chunks = []
            for i in range(0, len(text), chunk_size - chunk_overlap):
                chunks.append(text[i:i + chunk_size])
            return chunks
        texts = custom_text_splitter(file_content, chunk_size=800, chunk_overlap=80)
        embeddings = hf.embed_documents(texts)
        merged_embedding = np.concatenate([np.array(embed) for embed in embeddings], axis=0)

        # Index into Elasticsearch
        # for i, text in enumerate(texts):
        #     print("~~~~~~~~~~~~~:" + embeddings[i] + "\n\n")
        #     # metadata = {"filename": file_name, "url": file_url, "chunk_id": i}
        #     merge_text += embeddings[i]

        doc_body = {
            "content": file_content,
            "embedding": merged_embedding.tolist(),
            "metadata": {"filename": file_name, "url": file_url}
            # "metadata": metadata
        }
        es.index(index="search-my_elastic-2", document=doc_body)

    print("Finished indexing!")

if __name__ == "__main__":
    arguments = sys.argv
    if len(arguments) > 1:
        main_indexing(arguments[1])
    else:
        print("You need to provide a path to the folder with documents to index as a command line argument")

# Download and save the model locally (run this part separately before indexing)
# def download_model():
#     model_name = "sentence-transformers/all-MiniLM-L6-v2"
#     model = AutoModel.from_pretrained(model_name)
#     tokenizer = AutoTokenizer.from_pretrained(model_name)

#     model.save_pretrained("./local_model")
#     tokenizer.save_pretrained("./local_model")

# if __name__ == "__main__":
#     download_model()
